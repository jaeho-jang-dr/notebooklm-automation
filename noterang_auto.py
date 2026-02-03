#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
노트랑 완전 자동화 - 연구 → 슬라이드 생성 → 다운로드 → PPTX 변환
"""
import json
import subprocess
import sys
import asyncio
from pathlib import Path
from datetime import datetime
import time

if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

# 설정
NLM_EXE = Path.home() / "AppData/Roaming/Python/Python313/Scripts/nlm.exe"
DOWNLOAD_DIR = Path("D:/Entertainments/DevEnvironment/notebooklm")
AUTH_DIR = Path.home() / ".notebooklm-mcp-cli"

def sync_auth():
    """인증 동기화"""
    root_auth = AUTH_DIR / "auth.json"
    profile_dir = AUTH_DIR / "profiles" / "default"

    if not root_auth.exists():
        return False

    with open(root_auth) as f:
        root_data = json.load(f)

    cookies_dict = root_data.get('cookies', {})
    cookies_list = [
        {"name": n, "value": v, "domain": ".google.com", "path": "/",
         "expires": -1, "httpOnly": False, "secure": True, "sameSite": "Lax"}
        for n, v in cookies_dict.items()
    ]

    profile_dir.mkdir(parents=True, exist_ok=True)

    with open(profile_dir / "cookies.json", "w") as f:
        json.dump(cookies_list, f)

    with open(profile_dir / "metadata.json", "w") as f:
        json.dump({
            "csrf_token": root_data.get("csrf_token", ""),
            "session_id": root_data.get("session_id", ""),
            "email": root_data.get("email", ""),
            "last_validated": datetime.now().isoformat()
        }, f)

    with open(profile_dir / "auth.json", "w") as f:
        json.dump(root_data, f)

    return True

def run_nlm(args, timeout=120):
    """nlm CLI 실행"""
    import os
    cmd = [str(NLM_EXE)] + args
    env = os.environ.copy()
    env['PYTHONIOENCODING'] = 'utf-8'

    try:
        result = subprocess.run(
            cmd, capture_output=True, timeout=timeout,
            env=env
        )
        stdout = result.stdout.decode('utf-8', errors='replace') if result.stdout else ''
        stderr = result.stderr.decode('utf-8', errors='replace') if result.stderr else ''
        return result.returncode == 0, stdout, stderr
    except Exception as e:
        return False, '', str(e)

def check_auth():
    """인증 확인"""
    sync_auth()
    success, stdout, _ = run_nlm(["login", "--check"])
    return success and stdout and "valid" in stdout.lower()

def create_notebook(title):
    """노트북 생성"""
    success, stdout, stderr = run_nlm(["notebook", "create", title])
    if success:
        try:
            data = json.loads(stdout)
            return data.get('id')
        except:
            pass
    # 에러 메시지에서 ID 추출 시도
    return None

def research_and_import(notebook_id, query, mode="fast"):
    """연구 실행 및 소스 가져오기"""
    print(f"  검색: {query}")

    # 연구 시작
    success, stdout, _ = run_nlm([
        "research", "start", query,
        "--notebook-id", notebook_id,
        "--mode", mode
    ])

    if not success:
        return False, 0

    # Task ID 추출
    task_id = None
    for line in stdout.split('\n'):
        if 'Task ID:' in line:
            task_id = line.split('Task ID:')[1].strip()
            break

    if not task_id:
        return False, 0

    # 완료 대기
    print(f"  대기 중...", end="", flush=True)
    for i in range(24):  # 최대 2분
        time.sleep(5)
        success, stdout, _ = run_nlm(["research", "status", notebook_id])
        if "completed" in stdout.lower():
            print(" 완료!")
            break
        print(".", end="", flush=True)

    # 소스 가져오기
    success, stdout, _ = run_nlm(["research", "import", notebook_id, task_id])

    # 가져온 소스 수 추출
    imported = 0
    if "Imported" in stdout:
        try:
            imported = int(stdout.split("Imported")[1].split("source")[0].strip())
        except:
            pass

    return True, imported

def create_slides(notebook_id, language="ko", focus=None):
    """슬라이드 생성"""
    args = ["slides", "create", notebook_id, "--language", language, "--confirm"]
    if focus:
        args.extend(["--focus", focus])

    success, stdout, _ = run_nlm(args, timeout=60)

    if not success:
        return None

    # Artifact ID 추출
    artifact_id = None
    for line in stdout.split('\n'):
        if 'Artifact ID:' in line:
            artifact_id = line.split('Artifact ID:')[1].strip()
            break

    return artifact_id

def wait_for_slides(notebook_id, timeout=300):
    """슬라이드 생성 완료 대기"""
    print(f"  생성 중...", end="", flush=True)
    start = time.time()

    while time.time() - start < timeout:
        time.sleep(10)
        success, stdout, _ = run_nlm(["studio", "status", notebook_id])

        if '"status": "completed"' in stdout:
            print(" 완료!")
            # Artifact ID 추출
            try:
                data = json.loads(stdout)
                for item in data:
                    if item.get('type') == 'slide_deck' and item.get('status') == 'completed':
                        return item.get('id')
            except:
                pass
            return True

        elapsed = int(time.time() - start)
        print(f"\r  생성 중... {elapsed}초", end="", flush=True)

    print(" 타임아웃")
    return None

async def download_via_browser(notebook_id, output_dir):
    """브라우저를 통한 다운로드 (CLI 버그 우회)"""
    from playwright.async_api import async_playwright

    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    user_data_dir = AUTH_DIR / "browser_profile"

    async with async_playwright() as p:
        context = await p.chromium.launch_persistent_context(
            user_data_dir=str(user_data_dir),
            headless=False,
            downloads_path=str(output_dir),
            accept_downloads=True,
            args=['--disable-blink-features=AutomationControlled'],
            viewport={'width': 1920, 'height': 1080},
        )

        page = context.pages[0] if context.pages else await context.new_page()

        notebook_url = f"https://notebooklm.google.com/notebook/{notebook_id}"
        try:
            await page.goto(notebook_url, wait_until='domcontentloaded', timeout=30000)
        except:
            pass

        await asyncio.sleep(8)

        # 메뉴를 통한 다운로드
        downloaded_path = None
        menu_btns = await page.query_selector_all('[aria-haspopup="menu"], button[aria-label*="more"]')

        for menu_btn in menu_btns[-10:]:
            try:
                await menu_btn.click(force=True)
                await asyncio.sleep(1)

                dl_item = await page.query_selector('[role="menuitem"]:has-text("다운로드"), [role="menuitem"]:has-text("Download")')
                if dl_item:
                    async with page.expect_download(timeout=30000) as download_info:
                        await dl_item.click()

                    download = await download_info.value
                    filename = f"slides_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
                    downloaded_path = output_dir / filename
                    await download.save_as(str(downloaded_path))
                    break
            except:
                await page.keyboard.press('Escape')

        await asyncio.sleep(2)
        await context.close()

        return downloaded_path

def pdf_to_pptx(pdf_path):
    """PDF를 PPTX로 변환"""
    import fitz
    from pptx import Presentation
    from pptx.util import Inches
    import io

    pdf_path = Path(pdf_path)
    output_path = pdf_path.with_suffix('.pptx')

    doc = fitz.open(pdf_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for page_num in range(len(doc)):
        page = doc[page_num]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        img_data = pix.tobytes("png")

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(img_data),
            Inches(0), Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )

    doc.close()
    prs.save(output_path)

    return output_path, len(prs.slides)

def run_full_automation(title, research_queries, focus=None, language="ko"):
    """전체 자동화 실행"""
    print("=" * 60)
    print(f"노트랑 자동화: {title}")
    print("=" * 60)

    # 1. 인증 확인
    print("\n[1/6] 인증 확인...")
    if not check_auth():
        print("  ❌ 인증 실패. notebooklm-mcp-auth 실행 필요")
        return None
    print("  ✓ 인증 유효")

    # 2. 노트북 생성 또는 기존 사용
    print(f"\n[2/6] 노트북 확인...")
    success, stdout, _ = run_nlm(["list", "notebooks"])

    notebook_id = None
    if success:
        try:
            notebooks = json.loads(stdout)
            for nb in notebooks:
                if nb.get('title') == title:
                    notebook_id = nb.get('id')
                    print(f"  기존 노트북 사용: {notebook_id[:8]}...")
                    break
        except:
            pass

    if not notebook_id:
        notebook_id = create_notebook(title)
        if notebook_id:
            print(f"  새 노트북 생성: {notebook_id[:8]}...")
        else:
            print("  ❌ 노트북 생성 실패")
            return None

    # 3. 연구 및 소스 추가
    print(f"\n[3/6] 연구 자료 수집...")
    total_sources = 0
    for query in research_queries:
        success, count = research_and_import(notebook_id, query)
        total_sources += count
    print(f"  총 {total_sources}개 소스 추가")

    # 4. 슬라이드 생성
    print(f"\n[4/6] 슬라이드 생성...")
    artifact_id = create_slides(notebook_id, language=language, focus=focus)
    if not artifact_id:
        print("  ⚠️ 슬라이드 생성 시작 실패")

    # 5. 생성 완료 대기 및 다운로드
    print(f"\n[5/6] 다운로드...")
    wait_for_slides(notebook_id, timeout=300)

    # 브라우저로 다운로드 (CLI 버그 우회)
    pdf_path = asyncio.run(download_via_browser(notebook_id, DOWNLOAD_DIR))

    if not pdf_path or not pdf_path.exists():
        print("  ❌ 다운로드 실패")
        return None
    print(f"  ✓ PDF: {pdf_path.name}")

    # 6. PPTX 변환
    print(f"\n[6/6] PPTX 변환...")
    pptx_path, slide_count = pdf_to_pptx(pdf_path)
    print(f"  ✓ PPTX: {pptx_path.name} ({slide_count}슬라이드)")

    print("\n" + "=" * 60)
    print("완료!")
    print(f"  PDF:  {pdf_path}")
    print(f"  PPTX: {pptx_path}")
    print("=" * 60)

    return {
        'notebook_id': notebook_id,
        'pdf': str(pdf_path),
        'pptx': str(pptx_path),
        'slides': slide_count
    }

if __name__ == "__main__":
    # 예시 실행
    result = run_full_automation(
        title="견관절회전근개 파열",
        research_queries=[
            "회전근개 파열 원인 병인",
            "회전근개 파열 수술 치료",
            "회전근개 파열 재활 운동",
        ],
        focus="병인, 치료방법, 재활법",
        language="ko"
    )

    if result:
        print(f"\n결과: {json.dumps(result, ensure_ascii=False, indent=2)}")
